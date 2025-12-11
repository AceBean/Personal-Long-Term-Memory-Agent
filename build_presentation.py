#!/usr/bin/env python3
"""Build a PowerPoint (.pptx) from slides_markdown.md for the Personal Long-Term Memory Agent presentation.

Requirements:
- python-pptx (pip install python-pptx)

Usage:
- python build_presentation.py --input slides_markdown.md --output "Personal_Long-Term_Memory_Agent_Presentation.pptx" --presenter "Your Name"

The script parses the markdown slides separated by '---' and creates slides with title and bullet content.
"""

import argparse
from pptx import Presentation
from pptx.util import Pt
import re

SLIDE_SEPARATOR = r"^---\s*$"


def parse_slides(md_text):
    # Split on lines that contain only '---'
    parts = re.split(SLIDE_SEPARATOR, md_text, flags=re.MULTILINE)
    slides = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        # First line that starts with 'Slide' or a heading is title
        lines = [line.strip() for line in part.splitlines() if line.strip()]
        if not lines:
            continue
        # If first non-empty line starts with '#', use that as title
        if lines[0].startswith('#'):
            title = lines[0].lstrip('#').strip()
            body_lines = lines[1:]
        else:
            # Otherwise, take first line as title
            title = lines[0]
            body_lines = lines[1:]
        slides.append((title, body_lines))
    return slides


def add_bullets(slide, lines):
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    first = True
    for line in lines:
        if line.startswith('- '):
            text = line[2:]
            p = body.add_paragraph()
            p.level = 0
            p.text = text
            p.font.size = Pt(18)
        elif line.startswith('  - '):
            text = line.strip()[2:]
            p = body.add_paragraph()
            p.level = 1
            p.text = text
            p.font.size = Pt(16)
        else:
            # plain line, add as paragraph
            p = body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)


def build_pptx(input_md, output_pptx, presenter_name=None):
    with open(input_md, 'r', encoding='utf-8') as f:
        md_text = f.read()

    slides_data = parse_slides(md_text)
    prs = Presentation()
    # Set default slide size if desired

    for idx, (title, body_lines) in enumerate(slides_data):
        if idx == 0:
            # title slide layout
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            if '{{PRESENTER_NAME}}' in title and presenter_name:
                title = title.replace('{{PRESENTER_NAME}}', presenter_name)
            # title placeholder
            try:
                slide.shapes.title.text = title
            except Exception:
                pass
            # add subtitle if any body lines
            if body_lines:
                try:
                    subtitle = slide.placeholders[1]
                    subtitle.text = '\n'.join(body_lines)
                except Exception:
                    pass
        else:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            try:
                slide.shapes.title.text = title
            except Exception:
                pass
            add_bullets(slide, body_lines)

    prs.save(output_pptx)
    print(f"Saved presentation to {output_pptx}")


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', '-i', default='slides_markdown.md')
    parser.add_argument('--output', '-o', default='Personal_Long-Term_Memory_Agent_Presentation.pptx')
    parser.add_argument('--presenter', '-p', default='Your Name')
    args = parser.parse_args()
    build_pptx(args.input, args.output, presenter_name=args.presenter)
